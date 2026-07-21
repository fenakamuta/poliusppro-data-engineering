"""Gera o deck da Aula 3 em cima do arquivo da Aula 2 (herda tema, master e layouts)."""

import pathlib
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = "/home/nakamuta/Documents/poliusppro/aula-data-engineering/slides-e-materiais/Apresentacao Aula 2.pptx"
OUT  = "/home/nakamuta/Documents/poliusppro/aula-data-engineering/slides-e-materiais/Apresentacao Aula 3.pptx"

# --- trilhas (cores Open Color, mesma familia da aula 2) ---
TRACK = {
    "OLIST":  "10834A",  # verde  (dataset)
    "DOCKER": "0C8599",  # ciano
    "BANCO":  "1971C2",  # azul   (SQL/Postgres)
    "n8n":    "D6336C",  # rosa
    "DASH":   "7048E8",  # violeta (Metabase)
}
# cores do codigo (Material dark, iguais a aula 2)
C_BASE, C_COMMENT, C_STRING, C_KEYWORD = "E0E0E0", "6A737D", "C3E88D", "82AAFF"
KEYWORDS = set("""SELECT FROM WHERE GROUP BY ORDER DESC ASC LIMIT AS COUNT CREATE TABLE
PRIMARY KEY NOT NULL UNIQUE INSERT INTO VALUES EXPLAIN INDEX ON AND OR TRUE FALSE
TEXT INTEGER NUMERIC BOOLEAN TIMESTAMP SERIAL REFERENCES DEFAULT DISTINCT
DOCKER COMPOSE UP PS PULL NPX
SERVICES IMAGE ENVIRONMENT PORTS VOLUMES DEPENDS_ON""".split())

LAY = {"TITLE": 0, "OBJECT": 1, "SECTION": 2, "TITLE_ONLY": 5}

prs = Presentation(BASE)

# --- apaga todos os slides do arquivo base ---
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def rgb(h):
    return RGBColor.from_string(h)


def _fmt_run(r, size=None, bold=None, color=None, font=None):
    if size is not None:
        r.font.size = Pt(size)
    if bold is not None:
        r.font.bold = bold
    if color is not None:
        r.font.color.rgb = rgb(color)
    if font is not None:
        r.font.name = font


def set_title(slide, text):
    ph = slide.shapes.title
    ph.text_frame.clear()
    r = ph.text_frame.paragraphs[0].add_run()
    r.text = text
    _fmt_run(r, size=34, bold=True)


def _sem_bullet(p):
    """Desliga o marcador de lista que o layout aplica ao parágrafo."""
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    for bu in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(bu)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))
    pPr.set("indent", "0")
    pPr.set("marL", "0")


def fill_body(ph, lines, size=20, align=PP_ALIGN.LEFT):
    tf = ph.text_frame
    tf.clear()
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        _fmt_run(r, size=size)
        _sem_bullet(p)


def add_tag(slide, label):
    color = TRACK[label]
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(8.35), Inches(0.22), Inches(1.40), Inches(0.28))
    sp.fill.solid()
    sp.fill.fore_color.rgb = rgb(color)
    sp.line.fill.background()
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    _fmt_run(r, size=10, bold=True, color="FFFFFF", font="Menlo")


def new_slide(layout_key):
    return prs.slides.add_slide(prs.slide_layouts[LAY[layout_key]])


def slide_title(title, subtitle_lines):
    s = new_slide("TITLE")
    set_title(s, title)
    # subtitle placeholder = idx 1
    sub = s.placeholders[1]
    fill_body(sub, subtitle_lines, size=20, align=PP_ALIGN.CENTER)
    return s


def slide_section(title, subtitle):
    s = new_slide("SECTION")
    set_title(s, title)
    body = s.placeholders[1]
    fill_body(body, [subtitle], size=20)
    return s


def slide_object(title, lines, tag=None):
    s = new_slide("OBJECT")
    set_title(s, title)
    fill_body(s.placeholders[1], lines, size=20)
    if tag:
        add_tag(s, tag)
    return s


def slide_hook(title, lines, tag=None):
    s = new_slide("TITLE_ONLY")
    set_title(s, title)
    box = s.shapes.add_textbox(Inches(0.50), Inches(1.55), Inches(9.00), Inches(3.40))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = ln
        _fmt_run(r, size=28)
    if tag:
        add_tag(s, tag)
    return s


ASSETS = pathlib.Path(__file__).parent / "assets-slides"


def slide_diagram(title, img, caption=None, tag=None):
    """Slide visual: título + diagrama (PNG) + uma linha de fecho."""
    from PIL import Image

    s = new_slide("TITLE_ONLY")
    set_title(s, title)
    path = ASSETS / img
    with Image.open(path) as im:
        prop = im.height / im.width
    max_w, max_h = 9.20, (3.32 if caption else 3.85)
    w = max_w if max_w * prop <= max_h else max_h / prop
    h = w * prop
    s.shapes.add_picture(str(path), Inches((10 - w) / 2), Inches(1.32 + (max_h - h) / 2),
                         width=Inches(w))
    if caption:
        box = s.shapes.add_textbox(Inches(0.50), Inches(4.85), Inches(9.00), Inches(0.55))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = caption
        _fmt_run(r, size=15)
        r.font.italic = True
        r.font.color.rgb = RGBColor(0x49, 0x50, 0x57)
    if tag:
        add_tag(s, tag)
    return s


def tokenize(line):
    """Devolve lista de (texto, cor) com realce simples de sintaxe."""
    # acha inicio de comentario respeitando strings
    instr = False
    ci = None
    i = 0
    while i < len(line):
        ch = line[i]
        if ch == "'":
            instr = not instr
        elif not instr and (line[i:i+2] == "--" or ch == "#"):
            ci = i
            break
        i += 1
    code = line if ci is None else line[:ci]
    comment = "" if ci is None else line[ci:]
    segs = []
    for tok in re.findall(r"'[^']*'|[A-Za-z_][A-Za-z_0-9]*|[^A-Za-z_']+", code):
        if tok.startswith("'"):
            segs.append((tok, C_STRING))
        elif tok.upper() in KEYWORDS:
            segs.append((tok, C_KEYWORD))
        else:
            segs.append((tok, C_BASE))
    if comment:
        segs.append((comment, C_COMMENT))
    return segs


def slide_code(title, code_lines, tag=None, size=15):
    s = new_slide("OBJECT")
    set_title(s, title)
    # remove o placeholder de corpo vazio (nao usado nos slides de codigo)
    ph = s.placeholders[1]
    ph._element.getparent().remove(ph._element)
    # janela escura
    win = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.60), Inches(1.55), Inches(8.80), Inches(3.40))
    win.fill.solid()
    win.fill.fore_color.rgb = rgb("1E1E2E")
    win.line.fill.background()
    win.shadow.inherit = False
    # 3 botoes
    for j, col in enumerate(("FF5F56", "FFBD2E", "27C93F")):
        ov = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(0.80 + 0.28 * j), Inches(1.73), Inches(0.18), Inches(0.18))
        ov.fill.solid()
        ov.fill.fore_color.rgb = rgb(col)
        ov.line.fill.background()
        ov.shadow.inherit = False
    # codigo
    box = s.shapes.add_textbox(Inches(0.95), Inches(2.10), Inches(8.10), Inches(2.70))
    tf = box.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, ln in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if not ln:
            r = p.add_run(); r.text = " "; _fmt_run(r, size=size, color=C_BASE, font="Menlo")
            continue
        for txt, col in tokenize(ln):
            r = p.add_run()
            r.text = txt
            _fmt_run(r, size=size, color=col, font="Menlo")
    if tag:
        add_tag(s, tag)
    return s


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ============================ CONTEUDO ============================

# ---------------- ABERTURA ----------------
s = slide_title("Data Engineering III: do dado ao produto",
    ["Do dado guardado ao dashboard que se atualiza sozinho",
     "Prof. Fernando Henrique Vilella Nakamuta", "21/07/2026"])
note(s, "Boas-vindas. Hoje fechamos o ciclo: o dado sai da API, passa pelo banco e aparece num dashboard — no fim eles tem um SERVICO no ar, nao uma analise. AO LER O SUBTITULO, acrescente: 'e no fim eu conto qual e a pegadinha do se atualiza sozinho'. Isso evita prometer na abertura algo que o Bloco 3 vai qualificar.")

s = slide_hook("Lembram de onde paramos?",
    ["Na Aula 1 vocês exploraram o Olist:",
     "pedido atrasado vira review ruim (2.4 vs 4.3).",
     "",
     "Na Aula 2 treinaram um modelo (86%) e",
     "agendaram um workflow que coleta sozinho.",
     "",
     "A pergunta em que o notebook parou: onde o dado aterrissa?"])
note(s, "Recap das duas aulas em 30 segundos — e a PONTE com o notebook de revisão que acabou de rodar: ele terminou exatamente nesta pergunta. Deixe-a no ar — não responda ainda.")

s = slide_object("Os termos de hoje",
    ["Não precisa decorar nada disto agora.",
     "São as palavras que vão aparecer hoje —",
     "cada uma é traduzida na hora:",
     "",
     "container · imagem · volume        (o Docker)",
     "tabela · schema · chave primária · índice  (o banco)",
     "constraint · transação             (as regras)",
     "dashboard · string de conexão      (a entrega)",
     "",
     "(definições completas no material de apoio)"])
note(s, "NAO ler termo a termo — 10 segundos, so para dar familiaridade. Cada palavra e traduzida no slide em que aparece. CORTAVEL: se estiver atrasado, pule direto.")

s = slide_object("Onde estamos",
    ["Aula 1: Dados na prática (Olist)  ✓",
     "Aula 2: Do dado ao modelo + coleta  ✓",
     "Aula 3: Banco de dados + Dashboard  ← hoje",
     "Aula 4: ML vs. Agente de IA",
     "",
     "Hoje: guardar o dado, automatizar a entrada",
     "e entregar um dashboard para quem decide."])
note(s, "Rápido — 20 segundos. Se já estiver atrasado, fale por cima e pule.")

s = slide_diagram("O que vamos construir hoje", "d1_mapa.png",
    caption="Guarde este mapa: cada pedaço vira um bloco da aula.")
note(s, "O MAPA DA AULA — 60 segundos aqui valem ouro. Leia da esquerda pra direita: a API tem os pedidos, o n8n coleta e entrega, o Postgres guarda, o Metabase mostra, alguem DECIDE. O reloginho em cima do n8n e o 'sozinho' do titulo. Volte neste mapa na abertura de cada bloco: 'estamos AQUI'.")

s = slide_hook("Pra que um banco de dados?",
    ["Se o DuckDB é rápido, pra que um servidor?", "", "Essa é a pergunta do dia.", "No fim do Bloco 1 ela se responde sozinha."], tag="BANCO")
note(s, "A provocação central. NÃO responda agora — os próximos quatro slides são a resposta, e o recap do bloco 1 amarra.")

# ---------------- BLOCO 1 ----------------
s = slide_section("Bloco 1: Por que banco de dados",
    "Onde o dado vai morar.")
note(s, "Início. Energia alta. Bloco mais pesado do dia — controle o relógio.")

s = slide_object("Num arquivo, isso não fecha",
    ["O workflow da Aula 2 coleta às 8h. Mas coleta PARA ONDE?",
     "Imagine esse dado num Parquet e pergunte:",
     "",
     "• o n8n escreve às 8h enquanto alguém lê o dashboard?",
     "• chegou 1 pedido novo — reescreve o arquivo inteiro?",
     "• três pessoas abrem o dashboard — três cópias?",
     "",
     "Arquivo é dado parado, com um dono."], tag="BANCO")
note(s, "Os três furos do arquivo como destino de pipeline. Faça as três perguntas em voz alta e espere a turma responder. Concreto, com o Olist.")

s = slide_diagram("Banco: dado vivo, muitos donos", "d2_arquivo_vs_banco.png",
    caption="Planilha por e-mail vs. Google Sheets: todo mundo na mesma base, ao vivo.",
    tag="BANCO")
note(s, "FALE por cima do diagrama: um banco e um servidor SEMPRE LIGADO — varios sistemas leem e escrevem ao mesmo tempo, cada um ve o dado atualizado na hora. A analogia e-mail vs Google Sheets e a que pega: todo mundo de operacoes ja viveu o inferno do 'planilha_v3_final_REVISADA.xlsx'. Feche apontando o lado direito: o n8n grava ENQUANTO o Metabase le — isso se chama concorrencia.")

s = slide_diagram("DuckDB ou Postgres?", "d7_olap_oltp.png",
    caption="Não é qual é melhor — é qual para qual problema.", tag="BANCO")
note(s, "Abra as siglas POR CIMA do diagrama: OLAP = processamento ANALITICO (DuckDB, Aula 1: um arquivo, um usuario, roda e morre — poucas perguntas varrendo muito dado). OLTP = processamento TRANSACIONAL (Postgres, hoje: servidor ligado, muitos usuarios — muitas operacoes pequenas e simultaneas). Callback a licao da Aula 1: a ferramenta certa para o problema. CORTAVEL: se o Bloco 1 estiver atrasado, resuma em uma frase (arquivo para analisar, servidor para operar).")

s = slide_object("“Na minha máquina funciona”",
    ["Para o banco de hoje rodar, a máquina precisa de:",
     "Postgres na versão certa, configurado certo,",
     "com as dependências certas — e o mesmo pro Metabase.",
     "",
     "Só que aqui tem Windows, Mac e Linux.",
     "Instalar na mão = 40 máquinas, 40 resultados diferentes.",
     "",
     "O software depende do ambiente onde roda.",
     "Esse é o problema que o Docker resolve."], tag="DOCKER")
note(s, "Problema ANTES da ferramenta. A frase 'na minha máquina funciona' é a mais cara da computação — quem já pediu um relatório pro TI e ouviu 'aqui abre normal' conhece. Deixe o problema doer um pouco antes de apresentar a solução no próximo slide.")

s = slide_diagram("Docker: o container de navio", "d5_navio.png",
    caption="Empacota o programa com TUDO que ele precisa numa caixa padrão — e qualquer máquina roda.",
    tag="DOCKER")
note(s, "A analogia e feita para ESTA turma: gente de operacoes e logistica. Antes do container, cada carga era um problema — saca, barril, caixote, cada navio um quebra-cabeca. O container padronizou a CAIXA: nao importa o que tem dentro, qualquer navio e guindaste levam. Foi uma das maiores revolucoes de produtividade da logistica — e o Docker fez o MESMO com software (por isso o nome e o logo: uma baleia carregando containers). Fechamento: por isso o setup de casa foi 'instale o Docker', nao 'instale o Postgres' — instala-se o transportador uma vez, e ele carrega qualquer caixa.")

s = slide_diagram("Imagem, container e volume", "d3_imagem_container_volume.png",
    caption="Imagem = o molde · container = a caixa rodando · volume = onde o dado sobrevive.",
    tag="DOCKER")
note(s, "Tres palavras, um desenho. Imagem: o molde da caixa, baixa uma vez (docker pull). Container: a caixa rodando, criada do molde — DESCARTAVEL: apagou, some o que tinha dentro. Volume: a pasta que fica FORA da caixa, onde o dado sobrevive ao docker compose down. Volume e o que o aluno mais erra em casa: derruba a stack e acha que o banco sumiu. Aponte a linha volumes: no proximo slide.")

s = slide_code("A stack: docker-compose.yml",
    ["services:",
     "  postgres:                 # o banco  <- lembre deste nome",
     "    image: postgres:16",
     "    environment:",
     "      POSTGRES_USER: aluno",
     "      POSTGRES_PASSWORD: aula3",
     "      POSTGRES_DB: olist",
     "    ports: ['5432:5432']    # porta_da_maquina:porta_do_container",
     "    volumes: ['postgres-data:/var/lib/postgresql/data']",
     "  metabase:                 # o dashboard",
     "    image: metabase/metabase",
     "    ports: ['3000:3000']",
     "volumes:                    # sem declarar aqui, o compose recusa",
     "  postgres-data:"], tag="DOCKER", size=11)
note(s, "Leia linha a linha. image = qual molde baixar. environment = usuário/senha/banco. ports = 'porta da máquina : porta do container'. volumes = onde o dado sobrevive. IMPORTANTE: circule o nome do serviço 'postgres' — ele volta como a pegadinha nº1 no Bloco 3. O arquivo do repositório é este mesmo, com três extras que omiti do slide por espaço: um healthcheck, um depends_on (o Metabase espera o Postgres ficar pronto) e um volume próprio do Metabase (metabase-data + MB_DB_FILE, para os dashboards sobreviverem ao down). Se alguém perguntar, é isso.")

s = slide_code("Subir a stack e falar com ela",
    ["# 1. na pasta 2026/aula-3-banco-e-dashboard do repo:",
     "docker compose up -d",
     "docker compose ps          # os dois 'running'?",
     "",
     "# 2. abre o terminal do proprio Postgres (psql):",
     "docker compose exec postgres psql -U aluno -d olist",
     "",
     "# Postgres -> localhost:5432  |  Metabase -> localhost:3000",
     "# host = a maquina onde o banco esta; aqui, a sua"], tag="DOCKER", size=13)
note(s, "Dois comandos, um momento so: sobe e conversa. Sem o psql eles nao sabem ONDE digitar o SQL dos proximos slides. Quem nao tiver o repo: git clone https://github.com/fenakamuta/poliusppro-data-engineering (ou Download ZIP no site). Quem preferir clique: DBeaver, TablePlus ou pgAdmin, com host localhost, porta 5432, aluno/aula3/olist. PLANO B: quem nao baixou as imagens deixa o docker compose pull rodando em segundo plano (banda de casa da conta) e acompanha assistindo; retoma quando terminar. Nao pare a aula.")

s = slide_object("Do Parquet para a tabela",
    ["O Olist de hoje vem com duas colunas a mais,",
     "preenchidas pelo modelo que vocês treinaram na Aula 2:",
     "",
     "risco_prob     o termômetro: prob. de review ruim (0 a 1)",
     "risco_review   o alarme: risco_prob ≥ 0.5 → agir",
     "",
     "Já vêm PRONTAS no Parquet (do Release do GitHub) —",
     "hoje não treinamos nem rodamos o modelo de novo.",
     "",
     "É a descoberta da Aula 1 virando ação: agir ANTES da nota."], tag="OLIST")
note(s, "Responda a dúvida antes que ela apareça: a predição é um presente já pronto, ninguém precisa refazer a Aula 2. TERMOMETRO vs ALARME: o modelo devolve probabilidade (risco_prob); o corte em 0.5 vira decisão de negocio (risco_review) — quem quiser ser mais agressivo baixa a régua. E amarre no negócio: atraso -> review ruim (2.4 vs 4.3) é a razão de as colunas existirem.")

s = slide_code("CREATE TABLE: declare antes",
    ["CREATE TABLE pedidos (",
     "    pedido_id       text     PRIMARY KEY,  -- identidade unica",
     "    estado          text     NOT NULL,     -- obrigatorio",
     "    categoria       text,",
     "    preco           numeric,",
     "    prazo_dias      integer,",
     "    delivered_late  boolean,               -- atrasou? (aula 1)",
     "    review_score    integer,               -- a nota do cliente",
     "    risco_prob      numeric,   -- prob. de review ruim (termometro)",
     "    risco_review    boolean    -- risco_prob >= 0.5   (alarme)",
     ");"], tag="BANCO", size=13)
note(s, "Cada coluna tem tipo declarado. PRIMARY KEY = identidade única. O CSV aceitaria qualquer coisa calado; o banco recusa na porta. Note delivered_late e review_score: são as colunas da descoberta da Aula 1 — elas voltam na consulta de daqui a dois slides.")

s = slide_diagram("Schema, tipos e constraints", "d9_porteiro.png",
    caption="Schema = as tabelas e regras que você declarou. Constraint = a regra que o porteiro não deixa violar.",
    tag="BANCO")
note(s, "Leia o desenho: tres linhas chegam, o porteiro (as constraints) julga na porta. Tipo: texto onde devia ser numero? recusado. NOT NULL: obrigatorio vazio? recusado. PRIMARY KEY: id repetido? recusado. O arquivo aceitava tudo isso CALADO. Callback correto a Aula 2 ('a API traz dado sujo, sempre trate') — agora existe alguem que te OBRIGA a cumprir.")

s = slide_code("Carregar o Parquet na tabela",
    ["# tudo isto ja esta pronto no repo:  python carga.py",
     "import pandas as pd, psycopg2        # psycopg2 = quem fala Postgres",
     "from psycopg2.extras import execute_values",
     "cols = ['pedido_id','estado','categoria','preco','prazo_dias',",
     "        'delivered_late','review_score','risco_prob','risco_review']",
     "",
     "df = pd.read_parquet('olist_com_risco.parquet')[cols]   # ordem!",
     "df = df.astype(object).where(pd.notnull(df), None)      # NaN -> NULL",
     "conn = psycopg2.connect(host='localhost', port=5432,",
     "        user='aluno', password='aula3', dbname='olist')",
     "with conn.cursor() as cur:            # cursor = por onde o SQL passa",
     "    sql = f'INSERT INTO pedidos ({\", \".join(cols)}) VALUES %s'",
     "    execute_values(cur, sql, df.values.tolist())",
     "conn.commit()                         # so agora esta gravado"], tag="BANCO", size=11)
note(s, "ESTE CODIGO E O carga.py DO REPO. NAO digite nada: explique as linhas na tela e mande a TURMA RODAR O SCRIPT PRONTO agora — python carga.py (ele baixa o Parquet do Release sozinho, ~4 MB, e carrega 96.470 linhas). SEM ESTA CARGA O RESTO DA AULA NAO TEM DADO — o EXPLAIN, o count ao vivo e o dashboard dependem dela. Se falhar em alguem (ambiente Python), a pessoa acompanha assistindo e resolve no intervalo — deixe o chat aberto para duvidas. psycopg2 e pyarrow ja estao no requirements.txt do curso. ARMADILHA REAL (testada): review_score vem vazio em parte dos pedidos, o pandas traz NaN, e o Postgres recusa NaN em coluna integer — a linha do where() converte NaN em NULL. E a licao de dado sujo da Aula 2 na pratica. A lista cols e usada duas vezes — ordenar o dataframe E nomear as colunas do INSERT. O commit no fim e a deixa para o slide de transacao.")

s = slide_code("O mesmo SQL da Aula 1",
    ["-- AULA 1 (DuckDB, direto no arquivo):",
     "SELECT delivered_late, round(avg(review_score), 2) AS score",
     "FROM 'data/olist.parquet'",
     "GROUP BY delivered_late;",
     "",
     "-- HOJE (Postgres, numa tabela):",
     "SELECT delivered_late, round(avg(review_score), 2) AS score",
     "FROM pedidos",
     "GROUP BY delivered_late;"], tag="BANCO", size=13)
note(s, "Esta é A query da Aula 1 — a que revelou atraso -> review ruim (2.4 vs 4.3). ATENÇÃO AO NÚMERO: nesta tabela (deduplicada por pedido) a tela vai imprimir 2.27 vs 4.29 — mesmo fenômeno, recorte um pouco diferente; diga isso em meia frase se alguém notar. Mostre que só a linha do FROM mudou. Mensagem: SQL é portátil; muda o motor, não a linguagem.")

s = slide_diagram("Índice: o índice remissivo", "d6_indice.png",
    caption="Como achar uma palavra num livro: ler tudo, ou ir direto pelo índice remissivo.",
    tag="BANCO")
note(s, "Sem indice, para achar 'SP' o banco le a tabela INTEIRA — como ler o livro todo atras de uma palavra. O indice e o indice remissivo do fim do livro: 'SP: paginas 3, 47, 112' — vai direto. Custa espaco e deixa a escrita um tiquinho mais lenta; quase sempre vale.")

s = slide_code("EXPLAIN: o plano do banco",
    ["EXPLAIN ANALYZE SELECT * FROM pedidos WHERE estado = 'SP';",
     "",
     "  Seq Scan on pedidos               <- le as 96.470 linhas",
     "",
     "CREATE INDEX idx_estado ON pedidos(estado);",
     "-- rode o EXPLAIN de novo:",
     "",
     "  Bitmap Index Scan on idx_estado    <- usa o indice"], tag="BANCO", size=13)
note(s, "TESTADO com o dado real: o plano novo aparece como 'Bitmap Heap Scan' com um 'Bitmap Index Scan on idx_estado' dentro — o nome exato varia, o que importa e que o Seq Scan SUMIU e o idx_estado aparece. Nao prometa relogio: em tabela deste tamanho o tempo pode nao cair visivelmente; o que sempre muda e o PLANO. O PROXIMO SLIDE desintimida o log — nao explique tudo aqui.")

s = slide_diagram("EXPLAIN sem medo", "d15_explain.png",
    caption="No log inteiro, procure UMA coisa: “Seq Scan” = leu tudo · o nome do seu índice = foi direto.",
    tag="BANCO")
note(s, "O slide-calmante: o log do EXPLAIN intimida, e este desenho da a regra de leitura em 5 segundos — ignore custos, tempos e condicoes; procure Seq Scan OU o nome do indice. Se perguntarem o que e o 'cost': 'a regua interna do banco para comparar estrategias — quanto menor, mais barato; nao e tempo'. E siga.")

s = slide_object("JOGO: Adivinhe o plano",
    ["Cinco queries. Antes de cada EXPLAIN, vote:",
     "índice ou tabela inteira (Seq Scan)?",
     "",
     "1. WHERE estado = 'SP'",
     "2. WHERE pedido_id = '…'    ← PRIMARY KEY",
     "3. WHERE preco > 1000",
     "4. WHERE estado = 'AC' AND preco > 1000",
     "5. SELECT count(*)          ← pegadinha",
     "",
     "(script pronto: sql/03_jogo_adivinhe_o_plano.sql)"], tag="BANCO")
note(s, "ATIVIDADE ELASTICA (5-15 min) — use para ganhar tempo, corte sem do se atrasar. GABARITO TESTADO no dado real: 1) indice (Bitmap Index Scan em idx_estado); 2) Index Scan puro na PRIMARY KEY — todo PRIMARY KEY ganha indice de graca; 3) Seq Scan — nao ha indice em preco; 4) usa o indice do ESTADO e filtra preco depois; 5) Seq Scan — count(*) sem WHERE le a tabela inteira, indice nao ajuda. Cada um digita o palpite no chat antes do EXPLAIN; quem errar menos leva a gloria. A mensagem que fica: o banco ESCOLHE o plano, e o EXPLAIN mostra a escolha.")

s = slide_object("DINÂMICA: vote no chat",
    ["Por que o dashboard não pode ler direto do Parquet?",
     "",
     "A)  Parquet é um formato lento demais",
     "B)  arquivo é cópia parada, com um dono — não aguenta",
     "      leitura e escrita ao mesmo tempo",
     "C)  SQL não funciona em cima de arquivo",
     "D)  Parquet não guarda colunas numéricas",
     "",
     "(30 segundos: digitem a letra no chat)"], tag="BANCO")
note(s, "DINAMICA ELASTICA (3-5 min) — corte sem do se atrasar. GABARITO: B. Por que as outras: A e falsa (Parquet e rapidissimo para analise — Aula 1); C e a melhor pegadinha, porque eles FIZERAM SQL em arquivo com o DuckDB — o problema nao e o SQL, e a concorrencia; D e falsa obvia. Se a maioria votar C, otimo: vale 2 minutos explicando que ler da certo, o problema e todo mundo ESCREVENDO e lendo junto. E o formato exato da prova.")

s = slide_hook("Recap do Bloco 1",
    ["A pergunta que abriu o bloco: pra que um banco?",
     "",
     "Porque análise congelada cabe num arquivo,",
     "mas operação viva, com muita gente mexendo",
     "ao mesmo tempo, precisa de um servidor.",
     "",
     "Falta: fazer o dado chegar sozinho. → Bloco 2"])
note(s, "Feche explicitamente a pergunta da abertura ('Pra que um banco, se o DuckDB é rápido?') — é a promessa sendo paga. Só então abra o gancho do Bloco 2.")

s = slide_section("Intervalo", "15 minutos")
note(s, "Pausa. Voltamos às 20h25.")

# ---------------- BLOCO 2 ----------------
s = slide_section("Bloco 2: n8n alimentando o Postgres",
    "O workflow da Aula 2 finalmente aterrissa.")
note(s, "Início. Recuperar energia com 'fazer'. Bloco mais folgado do dia.")

s = slide_hook("O workflow ganhou um destino",
    ["Na Aula 2, o n8n coletava de uma API",
     "e… mostrava o JSON na tela.",
     "",
     "Hoje a estrutura é a mesma, com dois upgrades:",
     "a fonte é uma API de pedidos (já no repo)",
     "e o destino é um banco de verdade."])
note(s, "Ponte da Aula 2 — mas seja honesto sobre o que muda: a ESTRUTURA do fluxo e a mesma que eles ja montaram; a fonte e nova (a API de pedidos, slide adiante) e o destino e novo (o Postgres).")

s = slide_object("O node Postgres no n8n",
    ["No n8n, o Postgres é só mais um bloco. Ele pede:",
     "",
     "• host: em que máquina o banco está (localhost = a sua)",
     "• porta: qual porta de entrada dessa máquina (5432)",
     "• usuário, senha e banco (aluno / aula3 / olist)",
     "• operação: Insert",
     "",
     "São os mesmos 5 dados do docker-compose.yml."], tag="n8n")
note(s, "Traduza host e porta — é jargão puro para quem vem de operações. Porta = a porta de entrada da casa; o número diz qual serviço atende.")

s = slide_code("Mapear o JSON da API na tabela",
    ["# node Edit Fields  ->  renomeia campo a campo:",
     "",
     "pedido_id     <-  {{ $json.id }}",
     "estado        <-  {{ $json.customer_state }}",
     "categoria     <-  {{ $json.category }}",
     "preco         <-  {{ $json.price }}",
     "prazo_dias    <-  {{ $json.delivery_days }}",
     "# {{ $json.x }} = 'pega o campo x do que veio antes'",
     "# node Postgres  ->  Insert na tabela pedidos:",
     "# nomes iguais aos da coluna = mapeia sozinho"], tag="n8n", size=13)
note(s, "Mostre na tela do n8n: e um formulario, campo a campo. O {{ $json.campo }} e a sintaxe do n8n para 'pega este campo do que veio antes'. Como o Edit Fields ja renomeia para o nome exato das colunas, o node Postgres mapeia sozinho (auto-map). E exatamente assim no n8n-workflow-aula3.json do repo.")

s = slide_object("Cuidado: rode o n8n local",
    ["Rode com npx n8n, como na Aula 2 —",
     "ele abre em localhost:5678.",
     "",
     "O n8n.cloud NÃO enxerga o localhost:5432",
     "da sua máquina. Quem estiver na nuvem trava aqui.",
     "",
     "Local fala com local."], tag="n8n")
note(s, "PEGADINHA — anuncie ANTES de todo mundo começar a montar, não depois do erro. Pergunte no chat quem está no n8n.cloud e resolva agora.")

s = slide_code("De onde vêm os pedidos?",
    ["# o papel do sistema da empresa (ERP / e-commerce)",
     "# esta na pasta api-pedidos do repositorio:",
     "cd api-pedidos",
     "pip install -r requirements.txt",
     "uvicorn main:app --port 8001",
     "# dica: abra um NOVO terminal (serao 3: psql, API, n8n)",
     "# cada chamada devolve 10 pedidos NOVOS:",
     "# GET http://localhost:8001/pedidos",
     "# [{'id': 'novo-a1b2c3...', 'customer_state': 'SP',",
     "#   'category': 'housewares', 'price': 89.9, ...}]"], tag="n8n", size=13)
note(s, "Na vida real esse papel e do ERP/e-commerce da empresa; aqui, uma API do repo que sorteia pedidos reais do Olist com id sempre NOVO — por isso da para executar o workflow quantas vezes quiser sem duplicar PRIMARY KEY. Suba ANTES de montar o workflow e abra localhost:8001/pedidos no navegador para todos verem o JSON.")

s = slide_diagram("O workflow, ponta a ponta", "d8_workflow.png",
    caption="O garçom pede em localhost:8001/pedidos. Perdeu o passo? Importe n8n-workflow-aula3.json do repo.",
    tag="n8n")
note(s, "Monte ao vivo — e IMPORTANTE: antes, suba a API de pedidos do repo (pasta api-pedidos: uvicorn main:app --port 8001). O HTTP Request chama GET http://localhost:8001/pedidos e recebe 10 pedidos novos por execucao, com id inedito (pode executar quantas vezes quiser, nunca duplica PK). Os 3 primeiros blocos vieram da Aula 2 — retome a analogia do garcom. O bloco Postgres e a novidade. Tem um workflow pronto para importar no repo: n8n-workflow-aula3.json (so criar a credencial do Postgres).")

s = slide_code("Por baixo do bloco: um INSERT",
    ["-- o que o node Postgres executa, para cada item",
     "-- que a API devolveu:",
     "",
     "INSERT INTO pedidos",
     "  (pedido_id, estado, categoria, preco, prazo_dias)",
     "VALUES",
     "  ($1, $2, $3, $4, $5);",
     "",
     "-- $1..$5 sao buracos que o n8n preenche a cada linha.",
     "-- Nada a ver com o {{ $json.campo }} do slide anterior:",
     "-- aquele e a sintaxe do n8n; este e o SQL do Postgres."], tag="n8n", size=13)
note(s, "Não existe mágica no n8n. Abrir a caixa-preta aqui é o mesmo golpe didático do Metabase no Bloco 3: a interface bonita sempre vira SQL por baixo.")

s = slide_hook("Ao vivo: o dado chegou sozinho",
    ["SELECT count(*) FROM pedidos;   →   antes",
     "",
     "▶ executa o workflow",
     "",
     "SELECT count(*) FROM pedidos;   →   subiu +10.",
     "",
     "Ninguém digitou um INSERT."])
note(s, "Momento ao vivo. Rode o count ANTES na frente deles, execute, rode de novo. Deixe o silêncio trabalhar — não explique por cima.")

s = slide_diagram("De script a serviço: agendar", "d12_script_servico.png",
    caption="No n8n é um menu; por baixo é cron: 0 8 * * * = minuto 0, hora 8, todos os dias.",
    tag="n8n")
note(s, "Mecanica: troque o Manual Trigger pelo Schedule Trigger no proprio fluxo. Cuidado com o enquadre: agendar eles JA fizeram na Aula 2 — o pulo do gato de HOJE e que o workflow agendado agora tem um DESTINO que acumula. Script vs servico e a ideia que levam para casa (esta no desenho). Mostre o cron da caption para desmistificar o 'agendamento'.")

s = slide_code("O banco se defende do lixo",
    ["-- e se a fonte trouxesse nulo ou duplicata? o banco recusa",
     "",
     "INSERT INTO pedidos (pedido_id, estado)",
     "VALUES (NULL, 'SP');      -- ERRO: viola NOT NULL",
     "",
     "INSERT INTO pedidos (pedido_id, estado)",
     "VALUES ('abc', 'SP');     -- ok da 1a vez",
     "                          -- ERRO na 2a: PRIMARY KEY duplicada"], tag="BANCO", size=13)
note(s, "Tente inserir lixo AO VIVO e mostre o Postgres recusando. Ver o erro acontecer vale mais que dez slides sobre qualidade de dado.")

s = slide_object("JOGO: Quebre o banco",
    ["Desafio: inventem um INSERT que passe",
     "lixo pela porta. O chat dita, o banco julga.",
     "",
     "Spoiler: alguém vai conseguir —",
     "e é aí que a aula fica boa.",
     "",
     "(script pronto: sql/04_jogo_quebre_o_banco.sql)"], tag="BANCO")
note(s, "ATIVIDADE ELASTICA (5-15 min) — use para ganhar tempo, corte sem do se atrasar. O chat manda as tentativas, voce digita. O banco recusa NULL, id repetido, texto em coluna numerica... ate alguem tentar um VALOR ABSURDO valido: preco = -50 PASSA (testado!), estado = 'XX' passa, prazo = 9999 passa. REVELACAO: o banco so defende as regras que voce DECLAROU — constraint nao e magia, e contrato. Feche mostrando: ALTER TABLE pedidos ADD CONSTRAINT preco_positivo CHECK (preco >= 0); e o mesmo INSERT falhando. DEIXE O CHECK CRIADO (nao atrapalha nada) ou remova com DROP CONSTRAINT. E o gancho perfeito para o slide de transacao.")

s = slide_diagram("Ou tudo, ou nada: transação", "d10_transacao.png",
    caption="As garantias têm nome: ACID — Atomicidade, Consistência, Isolamento, Durabilidade.",
    tag="BANCO")
note(s, "Uma transacao agrupa operacoes numa coisa so: ou todas valem, ou nenhuma. Se a luz cai no meio da carga, o banco nao fica pela metade — desfaz e volta ao inicio. Abra a sigla ACID (esta na caption) — sem isso e decoreba. O exemplo classico em codigo esta no proximo slide.")

s = slide_code("A transferência, em código",
    ["-- ilustracao classica: a tabela contas NAO existe no olist",
     "BEGIN;",
     "",
     "UPDATE contas SET saldo = saldo - 100 WHERE id = 1;",
     "UPDATE contas SET saldo = saldo + 100 WHERE id = 2;",
     "",
     "-- se a luz cair aqui, NENHUMA das duas vale",
     "",
     "COMMIT;    -- ou ROLLBACK, e as duas voltam atras"], tag="BANCO", size=13)
note(s, "O dinheiro nao pode sumir entre duas contas. E o mesmo commit que apareceu no carregamento do Parquet, no Bloco 1. AVISE A TURMA: a tabela contas e so ilustracao classica, NAO existe no banco olist — ninguem precisa digitar isto, senao toma erro de relacao inexistente.")

s = slide_hook("Recap do Bloco 2",
    ["O dado chega sozinho, no horário,",
     "e cai num banco que recusa lixo.",
     "",
     "Falta: alguém ver isso sem abrir",
     "um terminal. → Bloco 3"])
note(s, "Fecha com a promessa do dashboard.")

s = slide_section("Intervalo", "15 minutos")
note(s, "Pausa. Voltamos às 21h50.")

# ---------------- BLOCO 3 ----------------
s = slide_section("Bloco 3: Dashboard no Metabase",
    "Tirar o dado do terminal e pôr na frente de quem decide.")
note(s, "Início. Reta final — aqui está o entregável. Controle o relógio para sobrar tempo do gancho da Aula 4.")

s = slide_diagram("O que é o Metabase", "d14_geladeira_vitrine.png",
    caption="A vitrine não guarda comida: ela mostra a da geladeira — sempre ao vivo.",
    tag="DASH")
note(s, "O Postgres guarda o dado — mas ninguem de operacoes vai abrir um terminal e escrever SQL. O Metabase conecta no banco e deixa montar perguntas e graficos SEM codigo. A analogia geladeira/vitrine e a que mais gruda no deck. Reforce: o Metabase NAO guarda dado — ele le do Postgres, sempre ao vivo.")

s = slide_object("Conectar o Metabase no Postgres",
    ["Abra localhost:3000, crie uma conta local",
     "(pode ser e-mail de mentira, é só na sua máquina)",
     "e clique em adicionar banco de dados.",
     "",
     "Tipo: PostgreSQL.",
     "Usuário aluno, senha aula3, banco olist.",
     "",
     "São os mesmos dados do psql. Menos um…"], tag="DASH")
note(s, "Deixe o 'menos um' no ar de propósito — o host é a pegadinha do próximo slide.")

s = slide_diagram("Cuidado: o host NÃO é localhost", "d4_host_metabase.png",
    caption="Dentro do Docker, localhost é o próprio container. O host é o NOME do serviço: postgres.",
    tag="DASH")
note(s, "PEGADINHA Nº1 DA AULA — meia turma vai errar, por isso o aviso vem ANTES do erro. Leia o desenho: o Metabase roda DENTRO do Docker; para ele, localhost e ele mesmo (o laco vermelho). O host certo e postgres — o NOME do servico la no docker-compose.yml (lembra do Bloco 1?). Voce, FORA do Docker, continua usando localhost normalmente. Volte no slide do compose e aponte o nome do servico.")

s = slide_object("A primeira pergunta, sem SQL",
    ["No construtor visual: pedidos por estado.",
     "Escolha a tabela, agrupe por estado, conte.",
     "",
     "Um gráfico, zero linhas de código.",
     "",
     "Agora clique em ver o SQL que ele gerou…"], tag="DASH")
note(s, "Construtor visual primeiro. O 'clique para ver o SQL' prepara a revelação do próximo slide.")

s = slide_code("O SQL que o Metabase escreveu",
    ["SELECT estado, count(*) AS pedidos",
     "FROM pedidos",
     "GROUP BY estado",
     "ORDER BY pedidos DESC;",
     "",
     "-- Nenhuma novidade: e um SELECT com GROUP BY,",
     "-- a mesma estrutura que voces escreveram na mao no Bloco 1."], tag="DASH")
note(s, "A revelação: a interface bonita vira o SQL que eles já sabem escrever. Desmistifica a ferramenta e valoriza o que aprenderam no Bloco 1.")

s = slide_object("DINÂMICA: vote no chat",
    ["Na conexão do Metabase, por que host = postgres,",
     "e não localhost?",
     "",
     "A)  porque postgres é mais rápido que localhost",
     "B)  porque localhost não existe no Windows",
     "C)  dentro do Docker, localhost é o próprio container;",
     "      o nome do serviço é o endereço na rede interna",
     "D)  porque o Metabase não aceita endereços IP",
     "",
     "(30 segundos: digitem a letra no chat)"], tag="DASH")
note(s, "DINAMICA ELASTICA (3-5 min) — corte sem do se atrasar. GABARITO: C — e a pegadinha n1 da aula, recem-vivida na conexao: quem errou a conexao ha 10 minutos agora fixa o porque. A e B sao falsas diretas; D e falsa (aceita IP normalmente). Se muita gente errar, volte 30 segundos no diagrama do host (o laco vermelho). E o formato exato da prova.")

s = slide_object("O dashboard “Pedidos em risco”",
    ["Cinco cartões, uma tela:",
     "",
     "• total de pedidos na base   ← é este que se move ao vivo",
     "• total em risco hoje",
     "• risco por estado · risco por categoria",
     "• a lista dos piores — para o time atacar amanhã",
     "",
     "Lembre por que isso importa: atraso vira nota 2.4."], tag="DASH")
note(s, "Este é o entregável da aula. Amarre no negócio (a descoberta da Aula 1): não é um gráfico bonito, é uma lista de pedidos para alguém salvar amanhã de manhã. NUMEROS REAIS do dado carregado: 96.470 pedidos no total, 19.426 em risco (20,1%), e o estado com mais risco é SP (6.686).")

s = slide_code("O SQL por trás dos cartões",
    ["-- 1) total de pedidos  (e este que sobe ao vivo)",
     "SELECT count(*) FROM pedidos;",
     "",
     "-- 2) total em risco",
     "SELECT count(*) FROM pedidos WHERE risco_review;",
     "",
     "-- 3) risco por estado  (no cartao seguinte, troque por categoria)",
     "SELECT estado, count(*) AS risco FROM pedidos",
     "WHERE risco_review GROUP BY estado ORDER BY risco DESC;",
     "",
     "-- 4) os piores, ordenados pelo risco REAL do modelo",
     "SELECT pedido_id, estado, round(risco_prob, 2) AS risco",
     "FROM pedidos WHERE risco_review",
     "ORDER BY risco_prob DESC LIMIT 20;"], tag="DASH", size=12)
note(s, "Deixe este slide na tela enquanto eles montam os cartoes. Cada cartao e uma pergunta — e toda pergunta vira SQL. ATENCAO: o cartao 1 (total, sem WHERE) e o que voce vai usar no momento ao vivo — os pedidos que o n8n insere entram sem risco_review preenchido, entao os cartoes com WHERE risco_review NAO se movem.")

s = slide_object("Montando o dashboard",
    ["1.  + Novo → Consulta SQL → cole o 1º SQL → Salvar",
     "     (nome: Total de pedidos)",
     "2.  Salvar em → Novo dashboard: Pedidos em risco",
     "3.  Repita para os outros quatro cartões",
     "4.  Arraste, redimensione, Salvar",
     "",
     "Pronto: a tela de quem decide, sem uma linha de código",
     "além do SQL que vocês já sabiam escrever."], tag="DASH")
note(s, "FACA O 1o CARTAO VOCE, devagar, projetado — os cliques sao pequenos e o caminho de menu confunde. Do 2o em diante a turma acompanha no proprio ritmo — fique de olho no chat. Nomeie os cartoes com nomes de gente (Total de pedidos, Em risco hoje) — e o que aparece no dashboard. O auto-refresh vem dois slides adiante, no momento da honestidade.")

s = slide_hook("O momento da aula",
    ["▶ executa o workflow no n8n",
     "↻ volta ao dashboard e aperta F5",
     "",
     "o total sobe +10, na frente de todo mundo",
     "",
     "O dado saiu da API, passou pelo banco e",
     "apareceu no gráfico. Você não tocou em nada."], tag="DASH")
note(s, "NUNCA CORTE ISTO — e a aula inteira num gesto. Use o cartao TOTAL DE PEDIDOS, nao o de risco: os pedidos que chegam pela API entram sem a coluna risco_review preenchida (o modelo nao roda ao vivo), entao o contador de risco nao se move. Ajuste os numeros ao que aparecer na tela. Deixe a turma ver o numero mudar antes de falar qualquer coisa — e guarde a pergunta obvia (quem preve o risco desse pedido novo?) para o gancho da Aula 4.")

s = slide_object("O que é automático aqui?",
    ["O F5 foi apertado na mão. Então, o que roda sozinho?",
     "",
     "A COLETA: o Schedule Trigger põe dado novo no banco",
     "às 8h, com você dormindo.",
     "",
     "E o painel também pode: no Metabase, engrenagem →",
     "Auto-refresh → a cada 1 minuto. Aí ninguém aperta nada.",
     "",
     "E repare: o contador de RISCO não se moveu. Guarde isso."], tag="DASH")
note(s, "Honestidade intelectual: um aluno atento vai notar que o F5 é manual e sentir que 'se atualiza sozinho' foi exagero. Antecipe isso — e mostre o auto-refresh do Metabase ao vivo se der tempo.")

s = slide_object("E na nuvem? Supabase",
    ["O Supabase é o mesmo Postgres que você subiu —",
     "só que hospedado por alguém, na nuvem.",
     "",
     "Vocês já usaram um, sem saber: o jogo da Aula 2",
     "gravava as jogadas num Supabase.",
     "",
     "Muda a string de conexão. Não muda o SQL."], tag="BANCO")
note(s, "Callback forte: o jogo 'Você vs. a Máquina' da Aula 2 rodava em cima de um Supabase, com um CREATE TABLE igualzinho ao do Bloco 1. Candidato a corte se o tempo apertar.")

s = slide_code("A string de conexão",
    ["# hoje, no seu Docker:",
     "postgresql://aluno:aula3@localhost:5432/olist",
     "",
     "# no Supabase (nuvem):",
     "postgresql://postgres:SENHA@db.xxxx.supabase.co:5432/postgres",
     "",
     "# usuario:senha@host:porta/banco",
     "# muda a linha. O SQL de cima continua identico."], tag="BANCO", size=12)
note(s, "Decomponha a string: usuário, senha, host, porta, banco — os mesmos 5 dados que apareceram no psql e no Metabase. Migrar de local para nuvem é trocar UMA linha.")

s = slide_hook("Isso é um produto de dados",
    ["Amanhã às 8h, alguém de operações abre esse",
     "dashboard e decide o que fazer com os pedidos",
     "que estão prestes a virar uma nota 2.",
     "",
     "Você não entregou uma análise.",
     "Você entregou um serviço."])
note(s, "Coroação. A diferença entre análise (foto) e serviço (algo que continua rodando depois que você fecha o notebook).")

s = slide_hook("O que você fez hoje",
    ["Subiu um banco, declarou um schema, carregou",
     "dados reais, fez um pipeline gravar sozinho",
     "e entregou um dashboard para quem decide.",
     "",
     "Aula 1 descobriu. Aula 2 previu.",
     "Aula 3 entregou. Isso é engenharia de dados."])
note(s, "Coroação final juntando o fio das três aulas. Ponta a ponta, hoje.")

s = slide_object("Para praticar em casa",
    ["Suba a stack e monte um dashboard com uma",
     "pergunta SUA sobre o Olist.",
     "",
     "Salve o SQL num arquivo .sql e mande no grupo",
     "da turma — quem usa GitHub pode subir no seu.",
     "",
     "(compose, SQL e scripts: pasta da aula 3 do repo)"])
note(s, "Tarefa: dashboard com pergunta propria + o SQL salvo e enviado no grupo. Git nunca foi ensinado no curso — por isso o caminho padrao e o grupo; GitHub fica como opcao para quem ja usa.")

s = slide_object("O gancho da Aula 4",
    ["A coluna risco_review veio de um modelo treinado,",
     "com dado rotulado, lá na Aula 2.",
     "",
     "Mas o pedido que entrou hoje pelo n8n chegou SEM",
     "risco — o modelo roda em lote, depois, no dado pronto.",
     "Quem decide o risco dele, AGORA?",
     "",
     "E se uma IA olhasse o pedido — prazo, categoria,",
     "preço — e decidisse na hora, sem treinar nada?"])
note(s, "Gancho da Aula 4, pagando a deixa do slide 'Sendo honesto' (o contador de risco nao se moveu). CUIDADO COM A LOGICA: o pedido novo NAO tem comentario ainda — o agente decide olhando os DADOS do pedido (prazo, categoria, preco) na hora, sem treino. E quando o review em portugues chegar, ele tambem le — isso e a Aula 4. NAO corte este slide: e o que segura o fio do curso. Se o tempo apertar, corte o Supabase e as leituras.")

s = slide_object("Leituras complementares",
    ["PostgreSQL — Documentation (postgresql.org/docs)",
     "Metabase — Getting Started (metabase.com/docs)",
     "n8n — Postgres node (docs.n8n.io)",
     "Docker — Get Started (docs.docker.com/get-started)",
     "Supabase — Docs (supabase.com/docs)"])
note(s, "Material de acesso livre. Pode virar só material pós-aula se o tempo apertar.")

s = slide_title("Obrigado.",
    ["Até a Aula 4", "Fernando Henrique Vilella Nakamuta"])
note(s, "Energia positiva.")

prs.save(OUT)
print(f"OK  {len(prs.slides._sldIdLst)} slides  ->  {OUT}")
